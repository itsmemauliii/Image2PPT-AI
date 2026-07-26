from pptx import Presentation
from pptx.util import Pt

def create_ppt(image,text_blocks):

    prs=Presentation()

    slide_layout=prs.slide_layouts[6]

    slide=prs.slides.add_slide(slide_layout)

    for item in text_blocks:

        left=item["box"][0][0]
        top=item["box"][0][1]

        textbox=slide.shapes.add_textbox(
            Pt(left),
            Pt(top),
            Pt(300),
            Pt(50)
        )

        frame=textbox.text_frame

        run=frame.paragraphs[0].add_run()

        run.text=item["text"]

        run.font.size=Pt(18)

    output="output/presentation.pptx"

    prs.save(output)

    return output
