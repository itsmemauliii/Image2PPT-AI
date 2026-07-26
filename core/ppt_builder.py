from pptx import Presentation
from pptx.util import Pt
import os

OUTPUT="output"

os.makedirs(OUTPUT,exist_ok=True)

def build_ppt(image,text):

    prs=Presentation()

    slide=prs.slides.add_slide(
        prs.slide_layouts[6]
    )

    for item in text:

        x=item["box"][0][0]

        y=item["box"][0][1]

        box=slide.shapes.add_textbox(

            Pt(x),

            Pt(y),

            Pt(250),

            Pt(40)

        )

        tf=box.text_frame

        tf.text=item["text"]

    file=os.path.join(

        OUTPUT,

        "presentation.pptx"

    )

    prs.save(file)

    return file
